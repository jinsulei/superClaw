#!/usr/bin/env python3
"""Offline document operations used by the portable Hermes runtime."""

from __future__ import annotations

import argparse
import json
import sys
import xml.etree.ElementTree as ElementTree
import zipfile
from pathlib import Path


def ensure_bundled_packages() -> None:
    """Allow the portable base interpreter to use Hermes' bundled packages."""
    runtime_dir = Path(__file__).resolve().parent.parent
    package_dir = runtime_dir / "hermes-agent" / "Lib" / "site-packages"
    if package_dir.is_dir() and str(package_dir) not in sys.path:
        sys.path.insert(0, str(package_dir))


ensure_bundled_packages()


def document_kind(path: Path) -> str:
    ext = path.suffix.lower()
    return {
        ".xlsx": "excel",
        ".xls": "excel",
        ".docx": "word",
        ".pdf": "pdf",
        ".pptx": "presentation",
    }.get(ext, "unknown")


def preview_excel(path: Path) -> dict:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=False)
    sheets = []
    for sheet in workbook.worksheets[:12]:
        rows = []
        for row in sheet.iter_rows(max_row=80, max_col=24, values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(values):
                rows.append(values)
        sheets.append({"name": sheet.title, "rows": rows, "maxRow": sheet.max_row, "maxColumn": sheet.max_column})
    return {"kind": "excel", "sheets": sheets}


def preview_word(path: Path) -> dict:
    from docx import Document

    document = Document(path)
    paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    tables = []
    for table in document.tables[:20]:
        rows = [[cell.text for cell in row.cells] for row in table.rows[:80]]
        tables.append(rows)
    return {"kind": "word", "paragraphs": paragraphs[:500], "tables": tables}


def preview_pdf(path: Path) -> dict:
    from pypdf import PdfReader

    reader = PdfReader(path)
    pages = []
    for index, page in enumerate(reader.pages[:80]):
        pages.append({"page": index + 1, "text": (page.extract_text() or "")[:12000]})
    metadata = {str(key): str(value) for key, value in (reader.metadata or {}).items()}
    return {"kind": "pdf", "pages": pages, "pageCount": len(reader.pages), "metadata": metadata}


PRESENTATION_TEXT_TAG = "{http://schemas.openxmlformats.org/drawingml/2006/main}t"


def presentation_slide_names(archive: zipfile.ZipFile) -> list[str]:
    """Return slides in their natural order without requiring python-pptx."""
    def slide_number(name: str) -> int:
        try:
            return int(Path(name).stem.replace("slide", ""))
        except ValueError:
            return 0

    return sorted(
        (
            name for name in archive.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ),
        key=slide_number,
    )


def preview_presentation(path: Path) -> dict:
    slides = []
    with zipfile.ZipFile(path) as archive:
        for index, name in enumerate(presentation_slide_names(archive), start=1):
            root = ElementTree.fromstring(archive.read(name))
            text = "\n".join(node.text or "" for node in root.iter(PRESENTATION_TEXT_TAG) if (node.text or "").strip())
            slides.append({"slide": index, "text": text[:24000]})
    return {"kind": "presentation", "slides": slides, "slideCount": len(slides)}


def preview(path: Path) -> dict:
    kind = document_kind(path)
    if kind == "excel":
        data = preview_excel(path)
    elif kind == "word":
        data = preview_word(path)
    elif kind == "pdf":
        data = preview_pdf(path)
    elif kind == "presentation":
        data = preview_presentation(path)
    else:
        raise ValueError(f"Unsupported document type: {path.suffix}")
    data.update({"ok": True, "path": str(path), "fileName": path.name})
    return data


def replace_excel(source: Path, output: Path, find: str, replace: str) -> dict:
    from openpyxl import load_workbook

    workbook = load_workbook(source)
    changed = 0
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and find in cell.value:
                    cell.value = cell.value.replace(find, replace)
                    changed += 1
    workbook.save(output)
    return {"ok": True, "kind": "excel", "output": str(output), "changed": changed}


def clean_excel(source: Path, output: Path) -> dict:
    """Clear empty-cell residue without changing values, formulas, or merged content."""
    from openpyxl import load_workbook
    from openpyxl.cell.cell import MergedCell
    from openpyxl.utils import column_index_from_string

    workbook = load_workbook(source, data_only=False)
    total_removed_cells = 0
    total_removed_dimensions = 0
    sheets = []

    for sheet in workbook.worksheets:
        content_rows = set()
        content_columns = set()
        stale_cells = []
        cells = list(sheet._cells.items())
        merged_ranges = tuple(sheet.merged_cells.ranges)

        def is_merged(cell) -> bool:
            if isinstance(cell, MergedCell):
                return True
            return any(cell.coordinate in merged for merged in merged_ranges)

        for (row, column), cell in cells:
            value = cell.value
            has_content = value is not None and (not isinstance(value, str) or value.strip() != "")
            has_content = has_content or cell.comment is not None or cell.hyperlink is not None
            if has_content:
                content_rows.add(row)
                content_columns.add(column)
            elif not is_merged(cell):
                stale_cells.append((row, column))

        for key in stale_cells:
            del sheet._cells[key]
        total_removed_cells += len(stale_cells)

        for row in list(sheet.row_dimensions):
            if row not in content_rows:
                del sheet.row_dimensions[row]
                total_removed_dimensions += 1
        for column in list(sheet.column_dimensions):
            try:
                index = column_index_from_string(column.split(":", 1)[0])
            except ValueError:
                continue
            if index not in content_columns:
                del sheet.column_dimensions[column]
                total_removed_dimensions += 1

        sheets.append({
            "name": sheet.title,
            "contentRows": len(content_rows),
            "contentColumns": len(content_columns),
            "removedBlankCells": len(stale_cells),
        })

    workbook.save(output)
    return {
        "ok": True,
        "kind": "excel",
        "output": str(output),
        "removedBlankCells": total_removed_cells,
        "removedBlankDimensions": total_removed_dimensions,
        "sheets": sheets,
    }


def replace_word(source: Path, output: Path, find: str, replace: str) -> dict:
    from docx import Document

    document = Document(source)
    changed = 0

    def replace_paragraph(paragraph):
        nonlocal changed
        if find in paragraph.text:
            for run in paragraph.runs:
                if find in run.text:
                    run.text = run.text.replace(find, replace)
                    changed += 1

    for paragraph in document.paragraphs:
        replace_paragraph(paragraph)
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    replace_paragraph(paragraph)
    document.save(output)
    return {"ok": True, "kind": "word", "output": str(output), "changed": changed}


def replace_presentation(source: Path, output: Path, find: str, replace: str) -> dict:
    if not find:
        raise ValueError("--find is required for presentation replacement")
    changed = 0
    with zipfile.ZipFile(source, "r") as reader, zipfile.ZipFile(output, "w") as writer:
        slide_names = set(presentation_slide_names(reader))
        for item in reader.infolist():
            payload = reader.read(item.filename)
            if item.filename in slide_names:
                root = ElementTree.fromstring(payload)
                for node in root.iter(PRESENTATION_TEXT_TAG):
                    if node.text and find in node.text:
                        node.text = node.text.replace(find, replace)
                        changed += 1
                payload = ElementTree.tostring(root, encoding="utf-8", xml_declaration=True)
            writer.writestr(item, payload)
    return {"ok": True, "kind": "presentation", "output": str(output), "changed": changed}


def create_presentation(output: Path, text: str) -> dict:
    """Create a simple, editable PPTX with the one bundled python-pptx package."""
    from pptx import Presentation
    from pptx.util import Pt

    try:
        spec = json.loads(text) if text.strip() else {}
    except json.JSONDecodeError:
        spec = {"title": text.strip() or "Untitled presentation"}
    if not isinstance(spec, dict):
        raise ValueError("Presentation specification must be a JSON object")

    title = str(spec.get("title") or "Untitled presentation").strip()
    raw_slides = spec.get("slides")
    if not isinstance(raw_slides, list) or not raw_slides:
        raw_slides = [{"title": title, "body": spec.get("body") or ""}]

    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = str(spec.get("subtitle") or "Created offline with SuperClaw")

    for item in raw_slides:
        item = item if isinstance(item, dict) else {"body": str(item)}
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = str(item.get("title") or title)
        body = item.get("body") or item.get("bullets") or ""
        lines = [str(value) for value in body if str(value).strip()] if isinstance(body, list) else [line.strip() for line in str(body).splitlines() if line.strip()]
        if len(slide.placeholders) > 1:
            frame = slide.placeholders[1].text_frame
            frame.clear()
            for index, line in enumerate(lines or [""]):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                paragraph.text = line
                paragraph.font.size = Pt(20)
    presentation.save(output)
    return {"ok": True, "kind": "presentation", "output": str(output), "slideCount": len(presentation.slides)}


def watermark_pdf(source: Path, output: Path, text: str) -> dict:
    from io import BytesIO
    from pypdf import PdfReader, PdfWriter
    from reportlab.pdfgen import canvas

    reader = PdfReader(source)
    writer = PdfWriter()
    for page in reader.pages:
        width = float(page.mediabox.width)
        height = float(page.mediabox.height)
        overlay_buffer = BytesIO()
        overlay = canvas.Canvas(overlay_buffer, pagesize=(width, height))
        overlay.setFont("Helvetica", 10)
        overlay.drawString(36, 24, text)
        overlay.save()
        overlay_buffer.seek(0)
        page.merge_page(PdfReader(overlay_buffer).pages[0])
        writer.add_page(page)
    writer.add_metadata(reader.metadata or {})
    with output.open("wb") as handle:
        writer.write(handle)
    return {"ok": True, "kind": "pdf", "output": str(output), "pageCount": len(reader.pages), "watermark": text}


def main() -> int:
    parser = argparse.ArgumentParser(description="SuperClaw Hermes offline document tool")
    parser.add_argument("command", choices=["preview", "replace", "clean-excel", "watermark", "create-presentation"])
    parser.add_argument("input", type=Path, nargs="?")
    parser.add_argument("--output", type=Path)
    parser.add_argument("--find", default="")
    parser.add_argument("--replace", default="")
    parser.add_argument("--text", default="")
    args = parser.parse_args()

    try:
        if args.command == "create-presentation":
            if not args.output:
                raise ValueError("--output is required when creating a presentation")
            output = args.output.resolve()
            if output.suffix.lower() != ".pptx":
                raise ValueError("Presentation output must use the .pptx extension")
            output.parent.mkdir(parents=True, exist_ok=True)
            result = create_presentation(output, args.text)
        else:
            if not args.input:
                raise ValueError("input is required for this document operation")
            source = args.input.resolve(strict=True)
            if args.command == "preview":
                result = preview(source)
            else:
                if not args.output:
                    raise ValueError("--output is required for edits")
                output = args.output.resolve()
                output.parent.mkdir(parents=True, exist_ok=True)
                kind = document_kind(source)
                if args.command == "replace" and kind == "excel":
                    result = replace_excel(source, output, args.find, args.replace)
                elif args.command == "clean-excel" and kind == "excel":
                    result = clean_excel(source, output)
                elif args.command == "replace" and kind == "word":
                    result = replace_word(source, output, args.find, args.replace)
                elif args.command == "replace" and kind == "presentation":
                    result = replace_presentation(source, output, args.find, args.replace)
                elif args.command == "watermark" and kind == "pdf":
                    result = watermark_pdf(source, output, args.text)
                else:
                    raise ValueError("This command is not supported for the document type")
        print(json.dumps(result, ensure_ascii=False))
        return 0
    except Exception as error:
        print(json.dumps({"ok": False, "error": str(error)}, ensure_ascii=False), file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
